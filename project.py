# -*- coding: utf-8 -*-
coordinate=[]

import numpy as np
import os, io
import json
import requests

import cv2
from pptx.util import Inches,Mm, Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE,MSO_VERTICAL_ANCHOR
from pptx import Presentation
from google.cloud import vision
from PIL import Image

os.environ["GOOGLE_APPLICATION_CREDENTIALS"]="root-amulet-358418-f65fae9b5f90.json"

text_number = 0
result_number=0

def DeleteAllFiles(filePath):
    if os.path.exists(filePath):
        for file in os.scandir(filePath):
            os.remove(file.path)
        return 'Remove All File'

    else:
        return 'Directory Not Found'

def loadimg(url):
    response = requests.get(imgurl)
    responselist = []
    responselist = json.loads(response.text)

    download_url = responselist['body']['URL']
    download_image = requests.get(download_url)

    photo = open('/tmp/sample.jpg', 'wb')
    photo.write(download_image.content)
    photo.close

def morphology(image):
    global result_number
    k = 0
    blur = cv2.GaussianBlur(img, (5, 5), 0)
    imgray = cv2.cvtColor(blur, cv2.COLOR_BGR2GRAY)
    th1 = cv2.adaptiveThreshold(imgray, 255, cv2.ADAPTIVE_THRESH_GAUSSIAN_C, cv2.THRESH_BINARY, 15, 8)  # gaussian
    kernel = np.ones((7, 3), np.uint8)
    kernel2= np.ones((5, 5), np.uint8)
    erode = cv2.morphologyEx(th1, cv2.MORPH_ERODE, kernel, iterations=3)
    erode2 = cv2.morphologyEx(th1, cv2.MORPH_ERODE, kernel2, iterations=3)


    height, width = erode.shape


    for j in range(height):
        for i in range(width):
            if (erode[j, i] < 128):
                erode[j, i] = 0

            else:
                erode[j, i] = 255



    #titles = ['Original', 'Erode','Diliation', 'Opening', 'Close']
    #images = [image, erode,diliation,opening,close]

    #for i in range(5):
    #    plt.subplot(2, 3, i + 1), plt.imshow(images[i], 'gray')
    #    plt.title(titles[i])
    #    plt.xticks([]), plt.yticks([])
    #plt.show()

    kernel = cv2.getStructuringElement(cv2.MORPH_ELLIPSE, (5, 5))
    thresh = cv2.morphologyEx(erode, cv2.MORPH_CLOSE, kernel)

    imFlood = thresh.copy()
    h, w = thresh.shape[:2]
    mask = np.zeros((h + 2, w + 2), np.uint8)
    cv2.floodFill(imFlood, mask, (0, 0), 0)


    # Combine flood filled image with original objects
    imFlood[np.where(thresh == 0)] = 255


    # Invert output colors
    imFlood = ~imFlood

    contours, high =cv2.findContours(imFlood, cv2.RETR_TREE, cv2.CHAIN_APPROX_SIMPLE)
    cv2.drawContours(img, contours, -1, (0, 255, 0), 1)

    for contour in contours:

        [x, y, w, h] = cv2.boundingRect(contour)

        if w * h < 900:
            continue
        k = k + 1
        result_number = result_number + 1

        coordinate.append([f'{k}.jpg', [x, y, w, h]])
        print(coordinate)

        cv2.rectangle(img, (x, y), (x + w, y + h), (255, 0, 255), 2)



        cropped_image = erode2[y:y + h, x:x + w]
        original_image = imgray[y: y + h, x: x + w]
        resize = cv2.resize(cropped_image, dsize=(28,28), interpolation=cv2.INTER_AREA)
        cv2.imwrite(f'/tmp/original_result{k}.jpg', original_image)
        cv2.imwrite(f'/tmp/result{k}.jpg', resize)



def image_crop(infilename, save_path):
        img = Image.open(infilename)
        (img_h, img_w) = img.size


        # crop ??? ????????? : grid_w, grid_h
        grid_w = 9.3  # crop width
        grid_h = 9.3  # crop height
        range_w = (int)(img_w / grid_w)
        range_h = (int)(img_h / grid_h)


        i = 0

        for w in range(range_w):
            for h in range(range_h):
                bbox = (h * grid_h, w * grid_w, (h + 1) * (grid_h), (w + 1) * (grid_w))

                # ?????? ?????? ??????, ?????? ?????? ???
                crop_img = img.crop(bbox)

                fname = f"{i}.jpg"
                savename = save_path + fname
                crop_img.save(savename)

                i += 1


def detect_text(path):
    from google.cloud import vision
    client = vision.ImageAnnotatorClient()

    with io.open(path, 'rb') as image_file:
        content = image_file.read()

    image = vision.Image(content=content)

    response = client.text_detection(image=image)
    texts = response.text_annotations

    for text in texts:
        return ('\n"{}"'.format(text.description))

def detect_text_fortext(path):
    from google.cloud import vision
    client = vision.ImageAnnotatorClient()

    with io.open(path, 'rb') as image_file:
        content = image_file.read()

    image = vision.Image(content=content)

    response = client.text_detection(image=image)
    texts = response.text_annotations
    w = 0
    global text_number
    for text in texts:

        textloc = ('\n"{}"'.format(text.description))

        verticesX = (['{}'.format(vertex.x) for vertex in text.bounding_poly.vertices])
        verticesY = (['{}'.format(vertex.y) for vertex in text.bounding_poly.vertices])
        xlist = [int(i) for i in verticesX]
        ylist = [int(i) for i in verticesY]
        xLength = xlist[1] - xlist[3]
        yLength = ylist[2] - ylist[0]
        w = w + 1
        text_number = text_number + 1
        print(text_number)
        textlist.append([f'{w}.jpg', [verticesX[3], verticesY[2], xLength, yLength], textloc])
        print(textlist)



def detect_figure():
    global result_number
    for f in range(result_number):
        im = cv2.imread(f'/tmp/result{f + 1}.jpg')

        imgray = cv2.cvtColor(im, cv2.COLOR_BGR2GRAY)

        ret, thresh = cv2.threshold(imgray, 80, 255, 1)  ### !! inverse binary !! ###

        contours, hierarchy = cv2.findContours(thresh, cv2.RETR_TREE, cv2.CHAIN_APPROX_SIMPLE)

        ctt = 0
        for cnt in contours:
            epsilon = 0.04 * cv2.arcLength(cnt, True)
            approx = cv2.approxPolyDP(cnt, epsilon, True)
            print(len(approx))  ### !! fix indentation !! ###
            if len(approx) == 5:
                print("pentagon")
                cv2.drawContours(im, [cnt], 0, 255, -1)
                coordinate[f].append("")
                coordinate[f].append("pentagon")
                break
            elif len(approx) == 3:
                print("triangle")
                cv2.drawContours(im, [cnt], 0, (0, 255, 0), -1)
                coordinate[f].append("")
                coordinate[f].append("triangle")
                break
            elif len(approx) == 4:
                print("square")
                cv2.drawContours(im, [cnt], 0, (0, 0, 255), -1)
                coordinate[f].append("")
                coordinate[f].append("square")
                break
            elif len(approx) >= 6:
                print("circle")
                cv2.drawContours(im, [cnt], 0, (0, 255, 255), -1)
                coordinate[f].append("")
                coordinate[f].append("circle")
                break
            else:
                print("unknown")
                cv2.drawContours(im, [cnt], 0, (0, 255, 255), -1)
                coordinate[f].append("")
                coordinate[f].append("unknown")
                break
        cv2.imwrite(f'/tmp/approx{f + 1}.jpg', im)

def makeppt():
    prs = Presentation()
    prs.slide_width = Cm(25)
    prs.slide_height = Cm(19)

    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    shapes = slide.shapes

    for i in range(2, int(text_number)+1):

        with open(f'/tmp/Txtdata{i}.json', 'r', encoding='UTF8') as file:  # k?????? Gvidata.json??? json_data??????????????????
            contents = file.read()  # string
            data = json.loads(contents)
            x = int (data["location"][0])
            y = int (data["location"][1])
            w = int (data["location"][2])
            h = int (data["location"][3])



            left = 250 * (x/W)
            top = 190 * (y/H)
            width = 250 * (w/W)
            height = 190 * (h/H)
            print(left, top, width, height)

            if data['text'] != "":
                print((data['text']).replace('"', ''))
                tb = slide.shapes.add_textbox(Mm(left), Mm(top), Mm(width), Mm(height))
                tf = tb.text_frame
                tf.text = ""
                tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
                p = tf.add_paragraph()
                p.text = (data['text']).replace('"', '')
                p.font.size = Pt(10)
                p.font.name = '?????? ??????'

    for e in range(2, result_number+1):
        with open(f'/tmp/Gvidata{e}.json', 'r', encoding='UTF8') as file:  # k?????? Gvidata.json??? json_data??????????????????
            contents = file.read()  # string
            data = json.loads(contents)
            x = int (data["location"][0])
            y = int (data["location"][1])
            w = int (data["location"][2])
            h = int (data["location"][3])



            left = 250 * (x/W)
            top = 190 * (y/H)
            width = 250 * (w/W)
            height = 190 * (h/H)


            if data['figure'] =="":
                continue
            else:
                if data['text'] != "":
                    print("pass text")
                if data['figure'] == "circle" and data['text'] == "":
                    print(data['figure'])
                    shape = shapes.add_shape(MSO_SHAPE.OVAL, Mm(left), Mm(top), Mm(width), Mm(height))
                    shape.fill.background()
                    line = shape.line
                    line.color.rgb = RGBColor(0, 0, 0)

                if data['figure'] == "square" and data['text'] == "":
                    print(data['figure'])
                    shape = shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Mm(left), Mm(top), Mm(width), Mm(height))
                    shape.fill.background()
                    line = shape.line
                    line.color.rgb = RGBColor(0, 0, 0)
                if data['figure'] == "triangle" and data['text'] == "":
                    print(data['figure'])
                    shape = shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Mm(left), Mm(top), Mm(width), Mm(height))
                    shape.fill.background()
                    line = shape.line
                    line.color.rgb = RGBColor(0, 0, 0)
                if data['figure'] == "unknown":
                    continue
    print("??????")
    prs.save('/tmp/demo.pptx')

    try:
        url = 'https://15zytiytli.execute-api.us-west-2.amazonaws.com/v2/uploadppt'
        files = {'file': open('/tmp/demo.pptx', 'rb')}
        r = requests.post(url, files=files)
        print(r.text)
    except:
        print("fail")

#?????? ??????


imgurl = 'https://15zytiytli.execute-api.us-west-2.amazonaws.com/v2/hknu-pptimage?file=image/8.jpg'

client = vision.ImageAnnotatorClient()
original_number = 0
textlist = []
a = []
xlist =[]
ylist =[]
f = 0

loadimg(imgurl)

img = cv2.imread('/tmp/sample.jpg')
H,W,C = img.shape

morphology(img)

for f in range(result_number):
    image_crop(f'/tmp/result{f + 1}.jpg', '/tmp/')
    print(f'/tmp/result/{f + 1}.jpg')

detect_figure()

dict_list = ['image_number','location','text','figure']
for k in range(result_number):
            dictionary = dict(zip(dict_list, coordinate[k]))
            print(json.dump(dictionary, open('/tmp/' + f'Gvidata{k + 1}.json', 'w')))



# detect_text and make text.json
txt_file_name =('/tmp/sample.jpg')
detect_text_fortext(txt_file_name)

dict_text_list = ['image_number','location','text']
print(text_number)
for l in range (text_number) :
    text_dictionary = dict(zip(dict_text_list,textlist[l]))
    print(json.dump(text_dictionary, open('/tmp/' + f'Txtdata{l + 1}.json', 'w')))



#make dummy text data to classify figure
for j in range(result_number):
    if j == 1:
        continue
    else:
        file_name =f'/tmp/original_result{j + 1}.jpg'
        #print(file_name)
        detect_text(file_name)

        if detect_text(file_name) == None:
            a.append('')
            original_number += 1
        else:
            a.append(detect_text(file_name))
        print(a)

for f in range(result_number):
    if f == 1:
        continue
    else:
        with open(f'/tmp/Gvidata{f + 1}.json', 'r', encoding='utf-8') as file:

            json_data = json.load(file)
            print(json_data)
            print(len(a))
            json_data['text'] = a[f-1]  # ?????? ????????? ?????? ????????????
            print(json.dump(json_data, open('/tmp/' + f'Gvidata{f + 1}.json', 'w')))  # ????????? json???????????? Gvidatak.json????????? ??????


makeppt()
'''
DeleteAllFiles('./sample_images/crop')
DeleteAllFiles('./sample_images/approx')
DeleteAllFiles('./sample_images/original_result')
DeleteAllFiles('./sample_images/result')
DeleteAllFiles('./json')
DeleteAllFiles('./text_json')
'''
exit()



