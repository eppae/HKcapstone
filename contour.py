import cv2
from matplotlib import pyplot as plt
img = cv2.imread('./sample_images/sample_ppt/7.jpg')
from PIL import Image
coordinate=[]
import matplotlib.pylab as plt
import cv2
import numpy as np
import glob
import os, io
import json
from pptx.util import Inches, Cm, Pt
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx import Presentation
from math import *
from google.cloud import vision

client = vision.ImageAnnotatorClient()
original_number = 0
textlist = []
a = []
H,W,C = img.shape
temp_len = len(os.listdir('./json/'))

def morphology(image):
    k = 0
    blur = cv2.GaussianBlur(img, (5, 5), 0)
    imgray = cv2.cvtColor(blur, cv2.COLOR_BGR2GRAY)
    th1 = cv2.adaptiveThreshold(imgray, 255, cv2.ADAPTIVE_THRESH_GAUSSIAN_C, cv2.THRESH_BINARY, 15, 8)  # gaussian
    kernel = np.ones((7, 3), np.uint8)
    kernel2= np.ones((5, 5), np.uint8)
    erode = cv2.morphologyEx(th1, cv2.MORPH_ERODE, kernel, iterations=3)
    erode2 = cv2.morphologyEx(th1, cv2.MORPH_ERODE, kernel2, iterations=3)


    height, width = erode.shape


    for j in range(height):
        for i in range(width):
            if (erode[j, i] < 128):
                erode[j, i] = 0

            else:
                erode[j, i] = 255



    #titles = ['Original', 'Erode','Diliation', 'Opening', 'Close']
    #images = [image, erode,diliation,opening,close]

    #for i in range(5):
    #    plt.subplot(2, 3, i + 1), plt.imshow(images[i], 'gray')
    #    plt.title(titles[i])
    #    plt.xticks([]), plt.yticks([])
    #plt.show()

    kernel = cv2.getStructuringElement(cv2.MORPH_ELLIPSE, (5, 5))
    thresh = cv2.morphologyEx(erode, cv2.MORPH_CLOSE, kernel)

    imFlood = thresh.copy()
    h, w = thresh.shape[:2]
    mask = np.zeros((h + 2, w + 2), np.uint8)
    cv2.floodFill(imFlood, mask, (0, 0), 0)


    # Combine flood filled image with original objects
    imFlood[np.where(thresh == 0)] = 255


    # Invert output colors
    imFlood = ~imFlood

    contours, high =cv2.findContours(imFlood, cv2.RETR_TREE, cv2.CHAIN_APPROX_SIMPLE)
    cv2.drawContours(img, contours, -1, (0, 255, 0), 1)

    for contour in contours:

        [x, y, w, h] = cv2.boundingRect(contour)

        if w * h < 900:
            continue
        k = k + 1


        coordinate.append([f'{k}.jpg', [x, y, w, h]])
        print(coordinate)

        cv2.rectangle(img, (x, y), (x + w, y + h), (255, 0, 255), 2)



        cropped_image = erode2[y:y + h, x:x + w]
        original_image = imgray[y: y + h, x: x + w]
        resize = cv2.resize(cropped_image, dsize=(28,28), interpolation=cv2.INTER_AREA)
        cv2.imwrite(f'./sample_images/original_result/{k}.jpg',original_image)
        cv2.imwrite(f'./sample_images/result/{k}.jpg', resize)



def image_crop(infilename, save_path):
        img = Image.open(infilename)
        (img_h, img_w) = img.size


        # crop 할 사이즈 : grid_w, grid_h
        grid_w = 9.3  # crop width
        grid_h = 9.3  # crop height
        range_w = (int)(img_w / grid_w)
        range_h = (int)(img_h / grid_h)


        i = 0

        for w in range(range_w):
            for h in range(range_h):
                bbox = (h * grid_h, w * grid_w, (h + 1) * (grid_h), (w + 1) * (grid_w))

                # 가로 세로 시작, 가로 세로 끝
                crop_img = img.crop(bbox)

                fname = f"{i}.jpg"
                savename = save_path + fname
                crop_img.save(savename)

                i += 1






def detect_text(path):
    from google.cloud import vision
    client = vision.ImageAnnotatorClient()

    with io.open(path, 'rb') as image_file:
        content = image_file.read()

    image = vision.Image(content=content)

    response = client.text_detection(image=image)
    texts = response.text_annotations

    for text in texts:
        return ('\n"{}"'.format(text.description))
        # print(textbox)

        # print(save_textlist)
        # with open('save_textlist.pkl','wb')as f:
        #    pickle.dump(save_textlist,f)

        # vertices = (['({},{})'.format(vertex.x, vertex.y)
        #           for vertex in text.bounding_poly.vertices])

        # print('bounds: {}'.format(','.join(vertices)))

        # with open("textdata.txt", 'a') as f:
        #    f.write(textbox)




def detect_figure():
    for f in range(len(os.listdir('./sample_images/result/'))):
        im = cv2.imread(f'./sample_images/result/{f + 1}.jpg')

        imgray = cv2.cvtColor(im, cv2.COLOR_BGR2GRAY)

        ret, thresh = cv2.threshold(imgray, 80, 255, 1)  ### !! inverse binary !! ###

        contours, hierarchy = cv2.findContours(thresh, cv2.RETR_TREE, cv2.CHAIN_APPROX_SIMPLE)

        ctt = 0
        for cnt in contours:
            epsilon = 0.04 * cv2.arcLength(cnt, True)
            approx = cv2.approxPolyDP(cnt, epsilon, True)
            print(len(approx))  ### !! fix indentation !! ###
            if len(approx) == 5:
                print("pentagon")
                cv2.drawContours(im, [cnt], 0, 255, -1)
                coordinate[f].append("")
                coordinate[f].append("pentagon")
                break
            elif len(approx) == 3:
                print("triangle")
                cv2.drawContours(im, [cnt], 0, (0, 255, 0), -1)
                coordinate[f].append("")
                coordinate[f].append("triangle")
                break
            elif len(approx) == 4:
                print("square")
                cv2.drawContours(im, [cnt], 0, (0, 0, 255), -1)
                coordinate[f].append("")
                coordinate[f].append("square")
                break
            elif len(approx) >= 6:
                print("circle")
                cv2.drawContours(im, [cnt], 0, (0, 255, 255), -1)
                coordinate[f].append("")
                coordinate[f].append("circle")
                break
            else:
                print("unknown")
                break
        cv2.imwrite(f'./sample_images/approx/{f + 1}.jpg', im)

def makeppt():
    prs = Presentation()
    prs.slide_width = Cm(25)
    prs.slide_height = Cm(19)

    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    shapes = slide.shapes

    for i in range(2, temp_len):

        with open( f'./json/Gvidata{i}.json', 'r', encoding='UTF8') as file:  # k번째 Gvidata.json을 json_data로읽어들인다
            contents = file.read()  # string
            data = json.loads(contents)
            x = int (data["location"][0])
            y = int (data["location"][1])
            w = int (data["location"][2])
            h = int (data["location"][3])



            left = round((25 * (x/W)),4)
            top = round((19 * ((y/H))),4)
            width = round((25 * ((w)/W)),4)
            height = round((19 * ((h)/H)),4)
            print(left, top, width, height)

            if data["text"] != "" and data["figure"] =="":
                print((data["text"]).replace('"',''))
                tb = slide.shapes.add_textbox(Cm(left), Cm(top), Cm(width), Cm(height))
                tf = tb.text_frame
                tf.text = ''
                p = tf.add_paragraph()
                p.text = (data["text"]).replace('"','')
                p.font.size = Pt(10)
            else:
                if data["figure"] == "circle":
                    print(data["figure"])
                    shape = shapes.add_shape(MSO_SHAPE.OVAL, Cm(left), Cm(top), Cm(width), Cm(height))
                    shape.fill.background()
                    line = shape.line
                    line.color.rgb = RGBColor(0, 0, 0)

                if data["figure"] == "rectangle":
                    print(data["figure"])
                    shape = shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(left), Cm(top), Cm(width), Cm(height))
                    shape.fill.background()
                    line = shape.line
                    line.color.rgb = RGBColor(0, 0, 0)
                if data["figure"] == "triangle":
                    print(data["figure"])
                    shape = shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Cm(left), Cm(top), Cm(width), Cm(height))
                    shape.fill.background()
                    line = shape.line
                    line.color.rgb = RGBColor(0, 0, 0)

                print("1")

    prs.save('demo.pptx')

#전체 로직
morphology(img)
f=0
for f in range(len(os.listdir('./sample_images/result/'))):
    image_crop(f'./sample_images/result/{f+1}.jpg', './sample_images/crop/')
    print(f'./sample_images/result/{f+1}.jpg')

detect_figure()
dict_list = ['image_number','location','text','figure']
for k in range(len(os.listdir('./sample_images/result/'))):
            dictionary = dict(zip(dict_list, coordinate[k]))
            print(json.dump(dictionary, open('./json/'+ f'Gvidata{k+1}.json', 'w')))


#detect_text
images = glob.glob('./sample_images/original_result/*.jpg')

for j in images:
    file_name = os.path.join('./sample_images/original_result', f'{j}')
    # print(file_name)
    detect_text(file_name)

    if detect_text(file_name) == None:
        a.append('')
        original_number += 1
    else:
        a.append(detect_text(file_name))

for k in range(2, len('./json/')):  # 1번 이미지는 추출할 이미지 전체 파일을 잡아서 2번부터, 시작 이미지는 1부터 시작해서 temp -1

    with open('./json/' + f'Gvidata{k}.json', 'r', encoding='UTF8') as file:  # k번째 Gvidata.json을 json_data로읽어들인다
        json_data = json.load(file)
        json_data['text'] = a[k - 1]  # 해당 텍스트 값을 집어넣기
        print(json.dump(json_data, open('./json/'+ f'Gvidata{k+1}.json', 'w')))  # 수정한 json데이터를 Gvidatak.json파일로 저장

#makeppt
makeppt()





