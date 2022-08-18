from math import *
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Cm, Pt

# Cm 단위
left = Cm(0)
top = Cm(0)
width = Cm(0)
height = Cm(0)

# Creating PowerPoint
prs = Presentation()
title_slide_layout = prs.slide_layouts[6]  # 빈 레이아웃
slide = prs.slides.add_slide(title_slide_layout)
shapes = slide.shapes

# text파일 열 경우
f = open('test.txt', 'r', encoding='UTF8')
t = f.read()
f.close()
print(t)

# text 임시 설정
left = top = width = height = Inches(1.5)

# add TextBox
textBox = shapes.add_textbox(left, top, width, height)
textFrame = textBox.text_frame

textFrame.text = t
#textBox.font.size = Pt(43)

# shape 임시 설정
label = "rectangle"
left = top = Inches(3)
width = Inches(5)
height = Inches(3)

# add Shape
if label == "circle":
    shape = shapes.add_shape(MSO_SHAPE.OVAL, left, top, width, height)
if label == "rectangle":
    shape = shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
if label == "triangle":
    shape = shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE,
                             left, top, width, height)
if label == "pentagon":
    shape = shapes.add_shape(MSO_SHAPE.REGULAR_PENTAGON,
                             left, top, width, height)

# arrow

prs.save('test.pptx')
