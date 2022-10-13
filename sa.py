import json
import os
from pptx.util import Inches, Cm, Pt
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx import Presentation
from math import *



prs = Presentation()
blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)
shapes = slide.shapes
for f in range(len(os.listdir('D:/practice/contour/json/'))):
    with open(f'D:/practice/contour/json/Gvidata{f+1}.json', 'r', encoding='utf-8') as f:
        contents = f.read()  # string
        data = json.loads(contents)
    print(data)
    print(data["location"])

    for i in range(0, len(data)):

        x = data["location"][0]
        y = data["location"][1]
        w = data["location"][2]
        h = data["location"][3]

        left = 25 * (int(x)/1433)
        top = 19*(int(y)/1037)
        width = 25 * ((int(w))/1433)
        height = 19 * ((int(h))/1037)
        print(left, top ,width, height)

    if data["figure"] == "":
            print(data["text"])

    else:
        if data["figure"] == "circle":
            print(data["figure"])
            shape = shapes.add_shape(MSO_SHAPE.OVAL, Cm(
                left), Cm(top), Cm(width), Cm(height))
            shape.fill.background()
            line = shape.line
            line.color.rgb = RGBColor(0, 0, 0)

        if data["figure"] == "rectangle":
            print(data["figure"])
            shape = shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(
                left), Cm(top), Cm(width), Cm(height))
            shape.fill.background()
            line = shape.line
            line.color.rgb = RGBColor(0, 0, 0)
        if data["figure"] == "triangle":
            print(data["figure"])
            shape = shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Cm(
                left), Cm(top), Cm(width), Cm(height))
            shape.fill.background()
            line = shape.line
            line.color.rgb = RGBColor(0, 0, 0)
        print("1")


prs.save('demo.pptx')

