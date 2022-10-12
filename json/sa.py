import json
from pptx.util import Inches, Cm, Pt
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx import Presentation
from math import *


with open('D:/practice/contour/json/Gvidata1.json', 'r', encoding='utf-8') as f:
    contents = f.read()  # string
    data = json.loads(contents)

prs = Presentation()
blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)
shapes = slide.shapes

print(data)
print(data["location"])
for i in range(0, len(data)):

    x = data["location"][0]
    y = data["location"][1]
    w = data["location"][2]
    h = data["location"][3]

    left = 24 * (int(x)/3365)
    top = 12 * (1 - int((h)/2386))
    width = 24 * ((int(w)-int(x))/3365)
    height = 12 * ((int(h)-int(y))/2386)
    print(x, y, w, h)

    if data["figure"] == "":
        print(data["text"])
        tb = slide.shapes.add_textbox(Cm(left), Cm(top), Cm(width), Cm(height))
        tf = tb.text_frame
        tf.text = ''
        p = tf.add_paragraph()
        p.text = data["text"]
        p.font.size = Pt(15)
    else:
        if data["figure"] == "circle":
            print(data["figure"])
            shape = shapes.add_shape(MSO_SHAPE.OVAL, Cm(
                left), Cm(top), Cm(width), Cm(height))
            shape.fill.background()
            line = shape.line
            line.color.rgb = RGBColor(0, 0, 0)
        if data["figure"] == "rectangle":
            print(data["figure"])
            shape = shapes.add_shape(MSO_SHAPE.OVAL, Cm(
                left), Cm(top), Cm(width), Cm(height))
            shape.fill.background()
            line = shape.line
            line.color.rgb = RGBColor(0, 0, 0)
        if data["figure"] == "triangle":
            print(data["figure"])
            shape = shapes.add_shape(MSO_SHAPE.OVAL, Cm(
                left), Cm(top), Cm(width), Cm(height))
            shape.fill.background()
            line = shape.line
            line.color.rgb = RGBColor(0, 0, 0)

prs.save('demo.pptx')